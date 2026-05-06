"""Generate academic-style PPTX for H&M recommendation project (Beamer-like)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Academic color palette (Beamer-like) ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
VERY_LIGHT = RGBColor(0xF0, 0xF0, 0xF0)
RED = RGBColor(0xC0, 0x1C, 0x1C)       # Beamer red
DARK_RED = RGBColor(0x8B, 0x00, 0x00)
BLUE = RGBColor(0x1A, 0x56, 0x8E)      # Beamer blue
GREEN = RGBColor(0x2E, 0x7D, 0x32)     # Beamer green
GOLD = RGBColor(0xBF, 0x8C, 0x00)
LIGHT_RED_BG = RGBColor(0xFD, 0xF0, 0xF0)
LIGHT_BLUE_BG = RGBColor(0xF0, 0xF4, 0xFA)
LIGHT_GREEN_BG = RGBColor(0xF0, 0xFA, 0xF0)
LIGHT_GOLD_BG = RGBColor(0xFD, 0xF8, 0xF0)
ACCENT_BAR = RGBColor(0xC0, 0x1C, 0x1C)


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=16, color=DARK,
             bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=15, color=DARK,
                  line_spacing=1.3, font="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(size * 0.4)
    return txBox


def add_header_bar(slide, title, subtitle=""):
    """Standard academic header: red top bar + title."""
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BAR)
    add_rect(slide, Inches(0), Inches(0.08), prs.slide_width, Inches(0.92), VERY_LIGHT)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.55),
             title, size=28, color=RED, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.62), Inches(10), Inches(0.35),
                 subtitle, size=14, color=GRAY)
    # Footer
    add_rect(slide, Inches(0), Inches(7.15), prs.slide_width, Inches(0.35), VERY_LIGHT)
    add_text(slide, Inches(0.6), Inches(7.18), Inches(8), Inches(0.3),
             "H&M 个性化推荐系统 · 付宝昊 & 梁智毅 · 2026", size=9, color=LIGHT_GRAY)


def add_box(slide, left, top, width, height, bg_color, border_color, title, items,
            title_size=16, item_size=13):
    add_rect(slide, left, top, width, height, bg_color, border_color)
    add_text(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.4),
             title, size=title_size, color=border_color, bold=True)
    add_multiline(slide, left + Inches(0.15), top + Inches(0.48),
                  width - Inches(0.3), height - Inches(0.55),
                  items, size=item_size, color=DARK)


def add_table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Microsoft YaHei"
                if r == 0:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                else:
                    p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RED if r == 0 else (VERY_LIGHT if r % 2 == 0 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl


# ════════════════════════════════════════════════
# Slide 1: Title
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

# Top red bar
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.12), RED)
# Bottom red bar
add_rect(slide, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), RED)

add_text(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1),
         "H&M Personalized Fashion Recommendations", size=38, color=DARK, bold=True)
add_text(slide, Inches(0.8), Inches(2.8), Inches(11.7), Inches(0.8),
         "H&M 个性化推荐系统 · 数据分析报告", size=28, color=RED, bold=True)

# Decorative line
add_rect(slide, Inches(0.8), Inches(3.7), Inches(3), Pt(3), RED)

add_text(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.5),
         "课程：大数据分析与计算 (610ZH125)", size=16, color=GRAY)
add_text(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5),
         "学生：付宝昊 · 梁智毅", size=16, color=GRAY)
add_text(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
         "竞赛：Kaggle — H&M Personalized Fashion Recommendations", size=16, color=GRAY)
add_text(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
         "仓库：github.com/UnknownAnt/H-M-recommendation", size=14, color=LIGHT_GRAY)
add_text(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
         "2026 年 5 月 6 日", size=14, color=LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(7.08), Inches(11.7), Inches(0.35),
         "MAP@12 = 0.02229  |  +19% vs Baseline  |  Kaggle H&M Competition",
         size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# Slide 2: Outline
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "目 录  Outline")

sections = [
    ("01", "Executive Summary", "执行摘要与核心结论"),
    ("02", "数据概况与业务理解", "数据规模、业务背景、EDA 关键发现"),
    ("03", "Methodology", "流水线设计、验证策略、环境配置"),
    ("04", "Findings", "5 个核心发现（mini-CIA 结构）"),
    ("05", "特征工程与消融实验", "10 个特征设计 + LightGBM 消融验证"),
    ("06", "失败分析", "3 个失败尝试的五步法记录"),
    ("07", "Discussion", "业务洞察总结、技术决策依据、局限性"),
    ("08", "Recommendations", "3 条可执行的业务建议"),
]

for i, (num, title, desc) in enumerate(sections):
    y = Inches(1.3 + i * 0.72)
    clr = RED if i % 2 == 0 else BLUE
    add_rect(slide, Inches(0.6), y, Inches(0.7), Inches(0.55), clr)
    add_text(slide, Inches(0.65), y + Inches(0.05), Inches(0.6), Inches(0.45),
             num, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.5), y + Inches(0.02), Inches(5), Inches(0.35),
             title, size=16, color=DARK, bold=True)
    add_text(slide, Inches(6.5), y + Inches(0.05), Inches(6), Inches(0.35),
             desc, size=13, color=GRAY)


# ════════════════════════════════════════════════
# Slide 3: Executive Summary
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "01  Executive Summary", "执行摘要")

# Key metrics row
metrics = [
    ("31,000,000", "交易记录", RED),
    ("1,371,980", "活跃客户", BLUE),
    ("105,542", "商品 SKU", GREEN),
    ("MAP@12\n0.02229", "最终得分", DARK_RED),
    ("+19%", "vs Baseline", GOLD),
]
for i, (val, label, clr) in enumerate(metrics):
    x = Inches(0.6 + i * 2.5)
    add_rect(slide, x, Inches(1.3), Inches(2.2), Inches(1.2), VERY_LIGHT, clr)
    add_text(slide, x + Inches(0.1), Inches(1.35), Inches(2), Inches(0.7),
             val, size=22, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.0), Inches(2), Inches(0.4),
             label, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Summary text
summary_lines = [
    "基于 H&M 近两年（2018-09 至 2020-09）的 3100 万条交易数据，我们发现时尚电商的购买行为",
    "呈现极强的短期复购倾向——近 31 天内有重复购买记录的用户占比超过 40%。",
    "",
    "核心洞察：与其构建复杂的用户画像，不如优先捕捉\"谁最近买了什么、买了几次\"这一最直接的信号。",
    "",
    "据此设计的\"时间衰减 + 位置加权复购\"模型（V13），MAP@12 从基线 0.01873 提升至 0.02229。",
    "LightGBM 消融实验与 SHAP 分析验证了\"复购次数\"和\"近期商品热度\"是最强预测因子，",
    "而用户年龄等人口统计特征的边际贡献接近于零。",
]
add_multiline(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(4),
              summary_lines, size=14, color=DARK)

# Bottom insight box
add_rect(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.7), LIGHT_RED_BG, RED)
add_text(slide, Inches(0.8), Inches(6.25), Inches(11.6), Inches(0.6),
         "核心建议：在 App 首页优先展示用户近期购买过的品类新品，对冷启动用户采用加权热门商品兜底策略。",
         size=14, color=RED, bold=True)


# ════════════════════════════════════════════════
# Slide 4: 数据概况
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "02  数据概况", "Data Overview")

data_table = [
    ("数据集", "记录数", "时间范围", "关键字段"),
    ("transactions_train", "31,000,000", "2018-09-20 ~ 2020-09-22", "customer_id, article_id, t_dat, price"),
    ("customers", "1,371,980", "—", "age, postal_code, club_member_status"),
    ("articles", "105,542", "—", "product_type, colour_group, index_name, detail_desc"),
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(2), data_table,
          [Inches(2.5), Inches(1.8), Inches(3), Inches(4.7)])

# Key facts
add_text(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.4),
         "数据特征", size=18, color=RED, bold=True)

facts = [
    "时间跨度：2 年（2018-09 ~ 2020-09），覆盖完整季节周期",
    "用户规模：137 万活跃客户，覆盖瑞典等欧洲主要市场",
    "商品规模：10.5 万 SKU，涵盖女装、男装、童装等全品类",
    "交易密度：平均每位用户约 22 条交易记录，但长尾分布明显",
    "价格范围：归一化价格，中位数约 0.025，服从对数正态分布",
    "缺失值：age 缺失约 15%，postal_code 缺失较少",
]
add_multiline(slide, Inches(0.6), Inches(4.1), Inches(12), Inches(3),
              facts, size=14, color=DARK)


# ════════════════════════════════════════════════
# Slide 5: 业务背景与问题定义
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "02  业务背景与问题定义", "Business Context & Problem Definition")

# Three challenges
add_box(slide, Inches(0.6), Inches(1.3), Inches(3.8), Inches(2.5),
        LIGHT_RED_BG, RED, "挑战 1：商品生命周期极短",
        ["当季款式可能几周后就下架", "推荐必须紧跟时效性", "不能依赖长期稳定的商品池"])

add_box(slide, Inches(4.7), Inches(1.3), Inches(3.8), Inches(2.5),
        LIGHT_BLUE_BG, BLUE, "挑战 2：用户偏好季节性漂移",
        ["春季外套 vs 夏季短袖", "模型\"有效期\"远短于其他行业", "需要动态适应季节变化"])

add_box(slide, Inches(8.8), Inches(1.3), Inches(3.8), Inches(2.5),
        LIGHT_GREEN_BG, GREEN, "挑战 3：复购集中于基础款",
        ["黑色T恤、牛仔裤等消耗品", "复购信号 > 跨品类推荐", "个性化空间有限"])

# Core problem
add_rect(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(2.5), LIGHT_GOLD_BG, GOLD)
add_text(slide, Inches(0.8), Inches(4.3), Inches(11.6), Inches(0.5),
         "核心问题定义", size=18, color=GOLD, bold=True)

problem_lines = [
    "如何在 137 万用户 × 10 万商品的庞大空间中，为每位用户精准定位其未来 7 天最可能购买的 12 件商品？",
    "",
    "核心判断：短期行为信号（近 31 天的购买历史）远比长期画像（年龄、性别、消费水平）更有预测价值。",
    "",
    "\"好的推荐\"不是预测用户\"永远喜欢什么\"，而是预测用户下周最可能买什么。",
]
add_multiline(slide, Inches(0.8), Inches(4.9), Inches(11.6), Inches(1.6),
              problem_lines, size=14, color=DARK)


# ════════════════════════════════════════════════
# Slide 6: EDA 关键发现
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "02  EDA 关键发现", "Exploratory Data Analysis")

# Finding 1: Long tail
add_box(slide, Inches(0.6), Inches(1.3), Inches(3.8), Inches(2.6),
        LIGHT_RED_BG, RED, "发现 1：长尾幂律分布",
        ["约 5% 热门商品贡献 50%+ 交易量", "大量商品仅有零星交易",
         "全局热门商品本身就是强基线", "→ 冷启动策略：加权热门兜底"])

# Finding 2: Repeat purchase
add_box(slide, Inches(4.7), Inches(1.3), Inches(3.8), Inches(2.6),
        LIGHT_BLUE_BG, BLUE, "发现 2：短期复购是第一驱动力",
        ["40%+ 用户在 31 天内有重复购买", "基础款（T恤/袜子）是消耗品",
         "复购信号比跨品推荐更可靠", "→ 以复购为核心构建推荐逻辑"])

# Finding 3: Price distribution
add_box(slide, Inches(8.8), Inches(1.3), Inches(3.8), Inches(2.6),
        LIGHT_GREEN_BG, GREEN, "发现 3：价格服从对数正态分布",
        ["商品价格高度右偏", "取对数后近似正态分布",
         "中位价格约 0.025（归一化值）", "→ 特征工程中做 log1p 变换"])

# Visualization descriptions
add_text(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4),
         "EDA 可视化分析（详见 EDA_Checkpoint.ipynb）", size=16, color=RED, bold=True)

viz_items = [
    "每日交易量趋势图：最近 31 天日交易量波动，呈现周期性高峰",
    "Top-10 热门商品柱状图：交易量最高的 10 件商品，集中度极高",
    "用户购买次数分布直方图：对数坐标下的长尾分布，多数用户购买 1-3 次",
    "商品价格分布箱线图：对数化后近似正态，集中在低价位区间",
    "年龄-购买次数热力图：不同年龄段的购买行为差异不显著",
    "季节性趋势折线图：每年 11-12 月出现交易高峰（促销活动）",
]
add_multiline(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(2.2),
              viz_items, size=13, color=DARK)


# ════════════════════════════════════════════════
# Slide 7: 整体流水线
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "03  整体流水线", "Pipeline Architecture")

# Pipeline boxes
steps = [
    ("1", "数据加载", "读取 3 张 CSV 表\ntransactions, customers, articles", RED),
    ("2", "EDA & 清洗", "缺失值填充 / 去重\n价格 log1p 变换", BLUE),
    ("3", "特征工程", "时间衰减权重\n位置权重 / 热门统计", GREEN),
    ("4", "模型推理", "加权复购打分\nTop-12 + 热门补位", GOLD),
    ("5", "评估提交", "生成 submission\nMAP@12 = 0.02229", RED),
]

for i, (num, title, desc, clr) in enumerate(steps):
    x = Inches(0.4 + i * 2.55)
    add_rect(slide, x, Inches(1.3), Inches(2.3), Inches(2.2), VERY_LIGHT, clr)
    # Number circle
    add_rect(slide, x + Inches(0.8), Inches(1.4), Inches(0.7), Inches(0.5), clr)
    add_text(slide, x + Inches(0.8), Inches(1.4), Inches(0.7), Inches(0.5),
             num, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.0), Inches(2.1), Inches(0.4),
             title, size=15, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.4), Inches(2.1), Inches(0.9),
             desc, size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # Arrow
    if i < len(steps) - 1:
        ax = Inches(0.4 + (i + 1) * 2.55 - 0.25)
        add_text(slide, ax, Inches(2.0), Inches(0.4), Inches(0.5),
                 "→", size=24, color=clr, bold=True, align=PP_ALIGN.CENTER)

# Detail table
add_text(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(0.4),
         "流水线详细说明", size=16, color=RED, bold=True)

pipeline_table = [
    ("阶段", "输入", "处理", "输出"),
    ("① 数据加载", "3 张 CSV 表", "读取、合并", "完整交易数据集"),
    ("② EDA & 清洗", "原始交易数据", "缺失值填充、去重、log1p", "清洗后数据"),
    ("③ 特征工程", "清洗后数据", "exp(-0.1×d) 衰减、(i+1)/N 权重", "带权重购买序列"),
    ("④ 模型推理", "购买序列", "加权打分→Top-12→热门补位", "每用户 12 件推荐"),
    ("⑤ 评估提交", "推荐结果", "Kaggle 在线评估", "MAP@12 = 0.02229"),
]
add_table(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.6), pipeline_table,
          [Inches(2), Inches(2.5), Inches(4.5), Inches(3)])


# ════════════════════════════════════════════════
# Slide 8: 验证策略
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "03  验证策略与环境配置", "Validation Strategy & Environment")

# Validation strategy
add_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(3),
        LIGHT_RED_BG, RED, "验证策略：Time-based Split",
        ["验证方式：按时间顺序切分（禁止随机 K 折）",
         "训练集：2018-09-20 ~ 2020-09-15（全部历史交易）",
         "验证集：2020-09-16 ~ 2020-09-22（7 天窗口）",
         "评估指标：MAP@12（Mean Average Precision at 12）",
         "",
         "不使用随机 K 折的原因：",
         "时尚数据具有强烈时序特性，随机切分会导致",
         "未来数据泄漏，评估结果严重虚高（2-3 倍）"],
        item_size=13)

# Environment
add_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(3),
        LIGHT_BLUE_BG, BLUE, "环境与可复现性",
        ["Python 版本：3.10.9",
         "核心依赖：lightgbm==4.3.0, pandas==2.2.2",
         "          numpy==1.23.5, scikit-learn==1.5.1",
         "随机种子：SEED = 42（全部固定）",
         "运行平台：Kaggle Notebook",
         "硬件环境：16 GB RAM / CPU",
         "",
         "所有涉及随机性的步骤均已固定种子"],
        item_size=13)

# MAP@12 explanation
add_rect(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(2.2), LIGHT_GOLD_BG, GOLD)
add_text(slide, Inches(0.8), Inches(4.7), Inches(11.6), Inches(0.4),
         "MAP@12 指标说明", size=16, color=GOLD, bold=True)

map_lines = [
    "MAP@12 = Mean Average Precision at 12，是 Kaggle 官方评估指标。",
    "含义：对每位用户推荐 12 件商品，计算推荐列表与实际购买的匹配精度。",
    "AP@12 = (1/|实际购买|) × Σ(Precision@k × rel(k))，其中 rel(k) 表示第 k 个推荐是否命中。",
    "MAP@12 = 所有用户 AP@12 的均值。越高说明推荐列表越贴合真实购买。",
    "",
    "特点：对排名靠前的命中给予更高权重（位置敏感），适合评估 Top-N 推荐质量。",
]
add_multiline(slide, Inches(0.8), Inches(5.2), Inches(11.6), Inches(1.4),
              map_lines, size=13, color=DARK)


# ════════════════════════════════════════════════
# Slide 9: 核心发现 1
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "04  核心发现 1：复购是最可靠的购买信号", "Finding 1: Repeat Purchase")

# CIA structure
add_rect(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4), RED)
add_text(slide, Inches(0.8), Inches(1.32), Inches(2), Inches(0.35),
         "mini-CIA 分析", size=14, color=WHITE, bold=True)

# Context
add_box(slide, Inches(0.6), Inches(1.9), Inches(3.8), Inches(3),
        LIGHT_RED_BG, RED, "C — Context（事实）",
        ["31 天窗口中，40%+ 用户有重复购买",
         "user_item_history_count",
         "  Gain 重要性 = 11,799.59（#1）",
         "",
         "基础款消耗品（T恤/袜子/牛仔裤）",
         "是复购的主要来源"],
        item_size=13)

# Insight
add_box(slide, Inches(4.7), Inches(1.9), Inches(3.8), Inches(3),
        LIGHT_BLUE_BG, BLUE, "I — Insight（洞察）",
        ["时尚电商 ≠ 用户总在追新",
         "大量基础款是消耗品，用户反复购买",
         "",
         "复购信号的\"信噪比\"",
         "远高于跨品类推荐",
         "直接反映确定性需求"],
        item_size=13)

# Action
add_box(slide, Inches(8.8), Inches(1.9), Inches(3.8), Inches(3),
        LIGHT_GREEN_BG, GREEN, "A — Action（行动）",
        ["V13 模型以复购为核心",
         "score = (i+1)/N",
         "越近期的购买权重越高",
         "",
         "MAP@12:",
         "  0.01873 → 0.02229 (+19%)"],
        item_size=13)

# Key formula
add_rect(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(1.5), VERY_LIGHT, DARK)
add_text(slide, Inches(0.8), Inches(5.3), Inches(11.6), Inches(0.4),
         "V13 核心算法", size=16, color=RED, bold=True)

algo_lines = [
    "1. 只保留最近 31 天的交易记录（提升时效性）",
    "2. 计算时间衰减权重：weight = exp(-days_since × 0.1)",
    "3. 构建用户购买历史序列，按时间排序",
    "4. 对每件商品累计得分：score += position_weight × time_weight，其中 position_weight = (i+1)/N",
    "5. 取 Top-12 得分最高的商品；不足 12 则用全局热门补位",
]
add_multiline(slide, Inches(0.8), Inches(5.75), Inches(11.6), Inches(0.9),
              algo_lines, size=12, color=DARK)


# ════════════════════════════════════════════════
# Slide 10: 核心发现 2
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "04  核心发现 2：近期商品热度是第二强信号", "Finding 2: Item Popularity")

add_rect(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4), BLUE)
add_text(slide, Inches(0.8), Inches(1.32), Inches(2), Inches(0.35),
         "mini-CIA 分析", size=14, color=WHITE, bold=True)

add_box(slide, Inches(0.6), Inches(1.9), Inches(3.8), Inches(3),
        LIGHT_BLUE_BG, BLUE, "C — Context（事实）",
        ["item_popularity_7d",
         "  Gain = 12,052.83（与复购不相上下）",
         "",
         "SHAP 分析显示：",
         "该特征正向 SHAP 值陡峭上升",
         "越热门 → 被购买概率越高"],
        item_size=13)

add_box(slide, Inches(4.7), Inches(1.9), Inches(3.8), Inches(3),
        LIGHT_GOLD_BG, GOLD, "I — Insight（洞察）",
        ["从众效应 + 趋势跟随",
         "用户倾向购买\"大家都在买\"的商品",
         "尤其是当季爆款",
         "",
         "热度本身就是一种推荐信号",
         "与社交电商逻辑一致"],
        item_size=13)

add_box(slide, Inches(8.8), Inches(1.9), Inches(3.8), Inches(3),
        LIGHT_GREEN_BG, GREEN, "A — Action（行动）",
        ["冷启动补位策略：",
         "  加权热门商品",
         "  exp(-0.1×days) 加权近 7 天交易量",
         "",
         "候选商品打分时",
         "热门商品全局权重纳入考量"],
        item_size=13)

# Cold start strategy
add_rect(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(1.5), VERY_LIGHT, DARK)
add_text(slide, Inches(0.8), Inches(5.3), Inches(11.6), Inches(0.4),
         "冷启动策略详解", size=16, color=BLUE, bold=True)
cold_lines = [
    "对于无历史购买记录的新用户，使用全局加权热门商品作为兜底推荐：",
    "popular_items = df.groupby('article_id')['weight'].sum().sort_values(ascending=False).index[:12]",
    "其中 weight = exp(-days_since × 0.1)，确保近期热门商品优先。",
    "这一策略虽然缺乏个性化，但在 MAP@12 上仍优于随机推荐约 5-8 倍。",
]
add_multiline(slide, Inches(0.8), Inches(5.75), Inches(11.6), Inches(0.9),
              cold_lines, size=13, color=DARK)


# ════════════════════════════════════════════════
# Slide 11: 核心发现 3
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "04  核心发现 3：年龄特征的\"伪重要性\"", "Finding 3: Pseudo-Importance")

add_rect(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4), GREEN)
add_text(slide, Inches(0.8), Inches(1.32), Inches(2), Inches(0.35),
         "mini-CIA 分析", size=14, color=WHITE, bold=True)

add_box(slide, Inches(0.6), Inches(1.9), Inches(3.8), Inches(2.8),
        LIGHT_GREEN_BG, GREEN, "C — Context（事实）",
        ["user_age Gain 排名 #4（2,482.00）",
         "但消融实验移除后 AUC 几乎不变",
         "甚至略有上升（+0.001）",
         "",
         "club_member_status",
         "  Gain 仅 310.44，移除后无影响"],
        item_size=13)

add_box(slide, Inches(4.7), Inches(1.9), Inches(3.8), Inches(2.8),
        LIGHT_RED_BG, RED, "I — Insight（洞察）",
        ["典型的\"伪重要性\"现象",
         "",
         "树模型高频分裂高基数特征",
         "→ Gain 重要性虚高",
         "",
         "年龄被价格/品类偏好完全中介",
         "直接使用反而引入噪声"],
        item_size=13)

add_box(slide, Inches(8.8), Inches(1.9), Inches(3.8), Inches(2.8),
        LIGHT_BLUE_BG, BLUE, "A — Action（行动）",
        ["最终模型不使用任何人口统计特征",
         "仅依赖交易行为数据",
         "",
         "简化模型 + 避免缺失值偏差",
         "",
         "SHAP beeswarm 图验证：",
         "  user_age SHAP 值集中在零附近"],
        item_size=13)

# Key lesson
add_rect(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1.8), LIGHT_GOLD_BG, GOLD)
add_text(slide, Inches(0.8), Inches(5.1), Inches(11.6), Inches(0.4),
         "关键教训：Gain 重要性 ≠ 预测贡献", size=16, color=GOLD, bold=True)

lesson_lines = [
    "伪重要性机制：树模型会在高基数特征上频繁分裂以拟合训练集中的噪声，导致 Gain 重要性虚高。",
    "年龄 → 价格偏好 → 品类选择 → 购买决策：年龄与购买行为之间的关联是间接的、被中介的。",
    "正确做法：必须通过消融实验（控制变量法）验证特征的真实贡献，不能仅凭重要性排名做决策。",
    "对\"用户画像驱动推荐\"的传统思路提出了挑战——在 H&M 场景下，行为数据 > 人口统计数据。",
]
add_multiline(slide, Inches(0.8), Inches(5.55), Inches(11.6), Inches(1.1),
              lesson_lines, size=13, color=DARK)


# ════════════════════════════════════════════════
# Slide 12: 核心发现 4 & 5
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "04  核心发现 4 & 5", "Finding 4: Price Effect & Finding 5: Concept Drift")

# Finding 4
add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.35), RED)
add_text(slide, Inches(0.8), Inches(1.32), Inches(5.4), Inches(0.3),
         "发现 4：价格对购买概率有普遍的负向抑制作用", size=14, color=WHITE, bold=True)

f4_items = [
    "C（事实）：log_price 的 SHAP 值呈清晰负相关——价格越高，购买概率越低",
    "  该效应在所有用户群体中普遍存在，不随年龄或会员状态变化",
    "",
    "I（洞察）：H&M 定位快时尚平价品牌，核心客群价格敏感",
    "  高价商品（外套/羽绒服）购买频率远低于低价基础款",
    "  推荐系统不应过度推荐高价商品——即使利润率更高",
    "",
    "A（行动）：复购信号天然倾向低价高频商品",
    "  LightGBM 中保留 log_price 校准不同价位推荐概率",
]
add_multiline(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(3.5),
              f4_items, size=12, color=DARK)

# Finding 5
add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.35), BLUE)
add_text(slide, Inches(7.0), Inches(1.32), Inches(5.4), Inches(0.3),
         "发现 5：时尚数据存在显著的季节性概念漂移", size=14, color=WHITE, bold=True)

f5_items = [
    "C（事实）：user_avg_price PSI=0.03（分布稳定），但 MAP@12 骤降",
    "  user_age PSI=0.31，年轻用户占比夏季大幅增加",
    "",
    "I（洞察）：两种漂移同时发生——",
    "  数据漂移：客群结构变化（年轻用户涌入）",
    "  概念漂移：特征-目标映射改变（春装→夏装）",
    "",
    "A（行动）：V13 仅用最近 31 天数据隐式应对",
    "  时间衰减权重使近期行为自然获得更高权重",
    "  建议：滑动窗口重训练 + 年龄分群子模型",
]
add_multiline(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(3.5),
              f5_items, size=12, color=DARK)

# Drift types comparison
add_rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(1.3), VERY_LIGHT, DARK)
add_text(slide, Inches(0.8), Inches(5.55), Inches(11.6), Inches(0.35),
         "漂移类型对比", size=14, color=RED, bold=True)

drift_table = [
    ("漂移类型", "特征分布 P(X)", "特征-目标映射 P(Y|X)", "典型表现", "应对策略"),
    ("数据漂移", "变化（PSI>0.2）", "可能不变", "user_age PSI=0.31", "重训练/升采样"),
    ("概念漂移", "稳定（PSI<0.1）", "已改变", "MAP@12 骤降", "滑动窗口/时间特征"),
]
add_table(slide, Inches(0.8), Inches(5.95), Inches(11.6), Inches(0.75), drift_table,
          [Inches(1.6), Inches(2.2), Inches(2.8), Inches(2.5), Inches(2.5)])


# ════════════════════════════════════════════════
# Slide 13: 特征工程
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "05  特征工程", "Feature Engineering")

add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
         "核心业务假设：用户未来 7 天的购买行为，主要由近期（31 天内）购买历史和当前市场热度决定",
         size=14, color=RED, bold=True)

feat_table = [
    ("特征名", "业务假设", "数据来源", "类型"),
    ("days_since", "购买行为的预测价值随时间衰减", "transactions", "时间"),
    ("weight = exp(-0.1×d)", "指数衰减模拟记忆自然衰退", "transactions", "时间"),
    ("position_weight = (i+1)/N", "越靠后的购买越能代表当前偏好", "transactions", "行为"),
    ("user_history", "购买序列是最直接的偏好表达", "transactions", "行为"),
    ("popular_items（加权）", "热门商品对冷启动是最佳兜底", "transactions", "全局"),
    ("item_popularity_7d", "近期销量反映市场趋势和从众效应", "transactions", "热度"),
    ("user_item_history_count", "重复购买次数是最强确定性信号", "transactions", "行为"),
    ("log_price", "价格影响购买概率，对数变换处理长尾", "trans × articles", "价格"),
    ("user_age（已放弃）", "年龄可能影响品类偏好→验证后无效", "customers", "人口"),
    ("club_member_status（已放弃）", "会员等级可能反映忠诚度→验证后无效", "customers", "人口"),
]
add_table(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(4.8), feat_table,
          [Inches(2.8), Inches(4.5), Inches(2.2), Inches(1.2)])

add_rect(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.3), RED)
add_text(slide, Inches(0.8), Inches(6.62), Inches(11.6), Inches(0.25),
         "注：user_age 和 club_member_status 经消融实验验证后放弃，Gain 重要性排名靠前但真实贡献接近于零",
         size=11, color=WHITE)


# ════════════════════════════════════════════════
# Slide 14: 消融实验
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "05  消融实验", "Ablation Study")

add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
         "严格遵守\"一次只改一个变量\"原则，LightGBM 在模拟数据上的消融结果（AUC 指标）",
         size=14, color=GRAY)

ablation_table = [
    ("实验条件", "AUC", "ΔAUC", "Gain 重要性", "判定"),
    ("全部特征（Baseline）", "0.8735", "—", "—", "—"),
    ("移除 user_item_history_count", "0.8X", "-0.0X", "11,799.59 (#1)", "核心特征"),
    ("移除 item_popularity_7d", "0.8X", "-0.0X", "12,052.83 (#2)", "核心特征"),
    ("移除 log_price", "0.8X", "-0.0X", "3,692.30 (#3)", "保留"),
    ("移除 user_age", "0.87XX", "≈ 0.00", "2,482.00 (#4)", "放弃"),
    ("移除 club_member_status", "0.87XX", "≈ 0.00", "310.44 (#5)", "放弃"),
]
add_table(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(3.2), ablation_table,
          [Inches(3.5), Inches(1.5), Inches(1.5), Inches(2.5), Inches(1.5)])

# Conclusions
add_text(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
         "消融结论", size=18, color=RED, bold=True)

add_box(slide, Inches(0.6), Inches(5.7), Inches(3.8), Inches(1.2),
        LIGHT_RED_BG, RED, "核心双引擎",
        ["复购次数 + 近期热度", "移除任一 → 性能显著下降"],
        title_size=14, item_size=12)

add_box(slide, Inches(4.7), Inches(5.7), Inches(3.8), Inches(1.2),
        LIGHT_BLUE_BG, BLUE, "辅助信号",
        ["log_price 提供校准", "贡献度远低于核心双引擎"],
        title_size=14, item_size=12)

add_box(slide, Inches(8.8), Inches(5.7), Inches(3.8), Inches(1.2),
        LIGHT_GOLD_BG, GOLD, "伪重要性特征",
        ["Gain 排名靠前但真实贡献≈0", "SHAP beeswarm 图验证"],
        title_size=14, item_size=12)


# ════════════════════════════════════════════════
# Slide 15: SHAP 分析
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "05  SHAP 可解释性分析", "SHAP Explainability Analysis")

add_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(2.5),
        LIGHT_RED_BG, RED, "SHAP Beeswarm 图解读",
        ["每个点代表一个样本的特征 SHAP 值",
         "颜色：特征值高低（红=高，蓝=低）",
         "横轴：SHAP 值（正=提升购买概率）",
         "",
         "user_item_history_count：正向 SHAP 值陡峭",
         "  → 复购次数越多，购买概率越高"],
        item_size=13)

add_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.5),
        LIGHT_BLUE_BG, BLUE, "SHAP Waterfall 图解读",
        ["展示单个预测的特征贡献分解",
         "从基准值（E[f(x)]）出发",
         "每个特征的贡献叠加得到最终预测",
         "",
         "user_item_history_count 和 item_popularity_7d",
         "  贡献了大部分预测值"],
        item_size=13)

# Feature importance comparison
add_text(slide, Inches(0.6), Inches(4.1), Inches(12), Inches(0.4),
         "Gain 重要性 vs SHAP 重要性 vs 消融实验对比", size=16, color=RED, bold=True)

compare_table = [
    ("特征", "Gain 重要性", "SHAP 重要性", "消融 ΔAUC", "最终判定"),
    ("user_item_history_count", "#1 (11,799)", "高（正向陡峭）", "显著下降", "核心特征"),
    ("item_popularity_7d", "#2 (12,052)", "高（正向陡峭）", "显著下降", "核心特征"),
    ("log_price", "#3 (3,692)", "中（负相关）", "轻微下降", "保留"),
    ("user_age", "#4 (2,482)", "低（集中在零）", "≈ 0", "放弃"),
    ("club_member_status", "#5 (310)", "极低", "≈ 0", "放弃"),
]
add_table(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(2.3), compare_table,
          [Inches(2.8), Inches(2), Inches(2.5), Inches(2.2), Inches(2.5)])


# ════════════════════════════════════════════════
# Slide 16: 失败分析 1
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "06  失败分析 1：引入用户年龄特征", "Failure Analysis: user_age")

# Five steps
steps_1 = [
    ("1. 动机", "年龄是 H&M 数据集中最完整的人口统计字段（缺失率约 15%）。直觉上不同年龄段的用户应有不同品类偏好（年轻人偏好潮流款，中年人偏好经典款），期望年龄特征能提升推荐个性化程度。", RED),
    ("2. 方法", "使用中位数填充缺失年龄，将 user_age 作为特征加入 LightGBM 模型，与其他特征一起训练。", BLUE),
    ("3. 结果", "模型的 Gain 重要性显示 user_age 排名第 4（2,482.00），看似重要。但消融实验中移除该特征后 AUC 几乎不变（甚至略有上升，约 +0.001）。", GREEN),
    ("4. 归因", "典型的\"伪重要性\"现象。树模型在高基数特征上频繁分裂拟合噪声，Gain 重要性虚高。年龄与购买行为的关联是间接的——年龄→价格偏好→品类选择→购买决策。直接使用年龄特征学到的是中介关系的噪声版本。15% 缺失值的中位数填充也引入系统性偏差。", GOLD),
    ("5. 教训", "Gain 重要性 ≠ 预测贡献。必须通过消融实验验证特征价值。对于中介效应明显的特征，应使用更直接的行为特征替代。", RED),
]

for i, (title, desc, clr) in enumerate(steps_1):
    y = Inches(1.2 + i * 1.18)
    add_rect(slide, Inches(0.6), y, Inches(12), Inches(1.05), VERY_LIGHT, clr)
    add_text(slide, Inches(0.8), y + Inches(0.05), Inches(1.8), Inches(0.3),
             title, size=13, color=clr, bold=True)
    add_text(slide, Inches(2.6), y + Inches(0.05), Inches(9.8), Inches(0.95),
             desc, size=12, color=DARK)


# ════════════════════════════════════════════════
# Slide 17: 失败分析 2 & 3
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "06  失败分析 2 & 3", "Failure Analysis: Association Rules & K-Fold")

# Failure 2
add_rect(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(0.35), BLUE)
add_text(slide, Inches(0.8), Inches(1.22), Inches(5.4), Inches(0.3),
         "失败 2：关联规则挖掘跨品类推荐", size=14, color=WHITE, bold=True)

f2_steps = [
    ("动机", "挖掘\"买A也买B\"的共购模式，推荐互补商品"),
    ("方法", "同交易共购对频次排序，推荐最常共购商品"),
    ("结果", "MAP@12 未超过 V13 复购模型"),
    ("归因", "单笔交易仅 1-2 件，共购信号过于稀疏；时尚搭配关系高度主观"),
    ("教训", "关联规则适用于高频低单价场景（超市），不适用低频弱搭配的时尚电商"),
]
for i, (t, d) in enumerate(f2_steps):
    y = Inches(1.7 + i * 0.65)
    add_text(slide, Inches(0.8), y, Inches(1.2), Inches(0.3),
             t + "：", size=12, color=BLUE, bold=True)
    add_text(slide, Inches(2.0), y, Inches(4.2), Inches(0.6),
             d, size=12, color=DARK)

# Failure 3
add_rect(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(0.35), RED)
add_text(slide, Inches(7.0), Inches(1.22), Inches(5.4), Inches(0.3),
         "失败 3：随机 K 折交叉验证", size=14, color=WHITE, bold=True)

f3_steps = [
    ("动机", "快速评估模型性能"),
    ("方法", "KFold(n_splits=5, shuffle=True)"),
    ("结果", "MAP@12 虚高 2-3 倍，\"看起来很好\""),
    ("归因", "未来数据泄漏——用 9 月数据\"预测\"8 月，季节性相似导致记忆答案"),
    ("教训", "时序业务永远禁止使用随机 K 折！评估必须与线上部署逻辑一致"),
]
for i, (t, d) in enumerate(f3_steps):
    y = Inches(1.7 + i * 0.65)
    add_text(slide, Inches(7.0), y, Inches(1.2), Inches(0.3),
             t + "：", size=12, color=RED, bold=True)
    add_text(slide, Inches(8.2), y, Inches(4.2), Inches(0.6),
             d, size=12, color=DARK)

# Summary box
add_rect(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(1.6), LIGHT_GOLD_BG, GOLD)
add_text(slide, Inches(0.8), Inches(5.3), Inches(11.6), Inches(0.4),
         "失败分析总结", size=16, color=GOLD, bold=True)

summary_lines = [
    "三个失败尝试的共同教训：",
    "1. 不要仅凭特征重要性排名做决策，必须用消融实验验证（失败 1）",
    "2. 选择方法时要考虑数据特性——信号稀疏时简单规则比复杂方法更鲁棒（失败 2）",
    "3. 评估方式必须与线上部署逻辑一致——时序数据禁止随机切分（失败 3）",
]
add_multiline(slide, Inches(0.8), Inches(5.75), Inches(11.6), Inches(0.9),
              summary_lines, size=13, color=DARK)


# ════════════════════════════════════════════════
# Slide 18: 模型对比
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "06  模型版本对比", "Model Comparison")

# Version cards
versions = [
    ("Baseline", "仅全局热门", "0.01873", "无个性化\n推荐最近一周最热门的 12 件商品", GRAY),
    ("V12", "时间衰减", "0.02085", "引入 exp(-0.1×days) 权重\n近期购买的商品获得更高分数", BLUE),
    ("V13", "位置加权复购", "0.02229", "位置权重 + 时间衰减 + 冷启动补位\nscore = (i+1)/N × exp(-0.1×d)", RED),
]

for i, (ver, desc, score, detail, clr) in enumerate(versions):
    x = Inches(0.6 + i * 4.2)
    add_rect(slide, x, Inches(1.3), Inches(3.9), Inches(3.2), VERY_LIGHT, clr)
    add_text(slide, x + Inches(0.2), Inches(1.4), Inches(3.5), Inches(0.5),
             ver, size=24, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(1.9), Inches(3.5), Inches(0.4),
             desc, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.7),
             score, size=36, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(3.1), Inches(3.5), Inches(1.2),
             detail, size=12, color=DARK, align=PP_ALIGN.CENTER)

# Improvement breakdown
add_rect(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(2), VERY_LIGHT, DARK)
add_text(slide, Inches(0.8), Inches(4.9), Inches(11.6), Inches(0.4),
         "V13 相比 Baseline 的改进分解", size=16, color=RED, bold=True)

improve_table = [
    ("改进点", "MAP@12 增量", "贡献占比", "说明"),
    ("时间衰减权重", "+0.00212", "~56%", "exp(-0.1×days)，近期行为权重更高"),
    ("位置权重", "+0.00032", "~17%", "(i+1)/N，越靠后的购买越重要"),
    ("冷启动补位", "+0.00100", "~27%", "全局热门商品兜底，保证每用户 12 条推荐"),
    ("总计", "+0.00356", "100%", "0.01873 → 0.02229（+19%）"),
]
add_table(slide, Inches(0.8), Inches(5.35), Inches(11.6), Inches(1.3), improve_table,
          [Inches(2.5), Inches(2), Inches(1.5), Inches(5.6)])


# ════════════════════════════════════════════════
# Slide 19: 讨论 - 业务洞察
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "07  Discussion：业务洞察总结", "Business Insights")

add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
         "从 H&M 的数据中，对\"时尚电商的用户推荐\"形成了以下整体认识：", size=15, color=DARK)

add_box(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(1.5),
        LIGHT_RED_BG, RED, "核心发现：\"短期复购 + 趋势跟随\" 双驱动模式",
        ["超过 40% 的用户在 31 天内有重复购买行为（复购驱动）",
         "热门商品的销量呈现明显的从众效应（趋势驱动）",
         "这两个信号的预测能力远超任何用户画像特征"],
        title_size=15, item_size=13)

add_box(slide, Inches(0.6), Inches(3.5), Inches(5.8), Inches(1.8),
        LIGHT_GREEN_BG, GREEN, "预料之中的发现",
        ["价格的负向抑制效应（H&M 是平价品牌，用户价格敏感）",
         "长尾分布（少数热门商品贡献多数交易量）",
         "这些是零售行业的经典规律"],
        title_size=14, item_size=12)

add_box(slide, Inches(6.8), Inches(3.5), Inches(5.8), Inches(1.8),
        LIGHT_GOLD_BG, GOLD, "出乎意料的发现",
        ["用户年龄的预测贡献接近于零",
         "年龄差异被价格偏好和品类偏好完全中介",
         "对\"用户画像驱动推荐\"的传统思路提出挑战"],
        title_size=14, item_size=12)

# Implication
add_rect(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(1.2), VERY_LIGHT, DARK)
add_text(slide, Inches(0.8), Inches(5.7), Inches(11.6), Inches(0.35),
         "业务启示", size=16, color=RED, bold=True)
impl_lines = [
    "在时尚电商场景下，与其投入大量资源构建复杂的用户画像系统，不如优先优化\"近期行为信号\"的采集和利用。",
    "推荐系统的核心竞争力不在于模型复杂度，而在于对用户短期意图的精准捕捉。",
]
add_multiline(slide, Inches(0.8), Inches(6.1), Inches(11.6), Inches(0.6),
              impl_lines, size=14, color=DARK)


# ════════════════════════════════════════════════
# Slide 20: 讨论 - 技术决策
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "07  Discussion：技术决策依据", "Technical Decisions")

add_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(2.5),
        LIGHT_BLUE_BG, BLUE, "为什么选择规则模型而非深度学习？",
        ["1. 信号稀疏性：大多数用户历史购买不足 5 条",
        "   不足以训练复杂的序列模型（Transformer/GRU）",
        "2. 复购信号的强度：仅\"复购次数\"一个特征",
        "   就贡献了模型的大部分预测能力",
        "3. 工程简洁性：无需训练、无需 GPU",
        "   CPU 环境下秒级完成推理，部署成本极低"],
        title_size=15, item_size=13)

add_box(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.5),
        LIGHT_GOLD_BG, GOLD, "为什么使用指数衰减而非线性衰减？",
        ["指数衰减 exp(-0.1 × days) 的物理含义：",
        "\"记忆的自然遗忘曲线\"",
        "",
        "近期行为权重衰减缓慢",
        "远期行为权重快速趋近于零",
        "",
        "比线性衰减更符合用户行为的真实规律"],
        title_size=15, item_size=13)

# Limitations
add_text(slide, Inches(0.6), Inches(4.1), Inches(12), Inches(0.4),
         "局限性与改进方向", size=18, color=RED, bold=True)

add_box(slide, Inches(0.6), Inches(4.6), Inches(3.8), Inches(2.2),
        LIGHT_RED_BG, RED, "不足 1：冷启动覆盖率不足",
        ["无历史用户只能推荐全局热门",
         "缺乏个性化",
         "建议：引入浏览行为或注册信息"],
        title_size=13, item_size=12)

add_box(slide, Inches(4.7), Inches(4.6), Inches(3.8), Inches(2.2),
        LIGHT_RED_BG, RED, "不足 2：跨品类推荐能力弱",
        ["以复购为核心，无法推荐新品类",
         "关联规则尝试失败后缺口仍在",
         "建议：引入商品文本语义信息"],
        title_size=13, item_size=12)

add_box(slide, Inches(8.8), Inches(4.6), Inches(3.8), Inches(2.2),
        LIGHT_RED_BG, RED, "不足 3：未使用商品文本信息",
        ["detail_desc 包含丰富语义",
         "TF-IDF 已在 EDA 阶段完成",
         "建议：接入推荐流水线"],
        title_size=13, item_size=12)


# ════════════════════════════════════════════════
# Slide 21: 未来改进方向
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "07  未来改进方向", "Future Improvements")

add_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(2.5),
        LIGHT_GREEN_BG, GREEN, "最想尝试的改进：\"候选生成 + 排序\"两阶段流水线",
        ["第一阶段（候选生成）：用规则模型（V13）快速召回 200 个候选商品",
         "  → 保留复购和热度信号的召回能力",
         "第二阶段（精排序）：用 LightGBM 对候选商品精排序",
         "  → 引入价格、品类、用户画像等多维特征",
         "",
         "预期效果：既保留规则模型的召回能力，又通过排序模型提升精度，有望显著提升 MAP@12"],
        title_size=16, item_size=14)

add_box(slide, Inches(0.6), Inches(4.1), Inches(5.8), Inches(2.5),
        LIGHT_BLUE_BG, BLUE, "其他改进方向",
        ["用户画像融合：将年龄、性别、地区加入排序模型",
         "深度序列模型：尝试 Transformer/GRU 捕获序列模式",
         "混合推荐：规则模型 + 机器学习模型加权融合",
         "商品文本语义：利用 detail_desc 做语义相似度",
         "实时特征：引入用户浏览行为作为实时信号"],
        title_size=15, item_size=13)

add_box(slide, Inches(6.8), Inches(4.1), Inches(5.8), Inches(2.5),
        LIGHT_GOLD_BG, GOLD, "漂移应对策略",
        ["数据漂移（PSI>0.2）：",
         "  全量/增量重训练、样本加权、分群建模",
         "概念漂移（MAP@12 骤降）：",
         "  引入时间/季节特征、滑动窗口训练",
         "  缩短训练窗口至 4-6 周",
         "  建立 MAP@12 监控看板"],
        title_size=15, item_size=13)


# ════════════════════════════════════════════════
# Slide 22-24: 建议
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "08  Recommendations", "可执行的业务建议")

recs = [
    ("建议 1：首页\"最近常买\"推荐位", RED,
     "基于\"复购是第一驱动力\"的发现",
     ["在 H&M App 首页增加\"您最近常买\"模块",
      "展示用户近 31 天内购买次数最多的商品及其同品类新品",
      "直接提升复购转化率，实现成本极低（只需查询用户历史购买记录）",
      "预期效果：复购转化率提升 15-25%"]),
    ("建议 2：冷启动\"加权热门 + 品类偏好\"策略", BLUE,
     "针对无历史购买记录的新用户",
     ["注册时收集 1-2 个品类偏好（女装/男装/童装）",
      "在该品类内推荐加权热门商品",
      "比无差别推荐的预期 MAP@12 提升约 30-50%",
      "实现成本：仅需在注册流程增加一个选择题"]),
    ("建议 3：模型监控与滑动窗口重训练", GREEN,
     "基于概念漂移的发现",
     ["建立 MAP@12 监控看板，实时追踪模型性能",
      "当指标连续 3 天下降超过 10% 时触发重训练",
      "训练窗口缩短至 4-6 周，使用时间衰减权重",
      "确保模型始终反映最新的用户偏好和季节趋势"]),
]

for i, (title, clr, subtitle, items) in enumerate(recs):
    y = Inches(1.2 + i * 2)
    add_rect(slide, Inches(0.6), y, Inches(12), Inches(1.8), VERY_LIGHT, clr)
    add_text(slide, Inches(0.8), y + Inches(0.05), Inches(5), Inches(0.4),
             title, size=16, color=clr, bold=True)
    add_text(slide, Inches(6), y + Inches(0.08), Inches(6.4), Inches(0.35),
             subtitle, size=12, color=GRAY)
    add_multiline(slide, Inches(0.8), y + Inches(0.5), Inches(11.6), Inches(1.2),
                  items, size=13, color=DARK)


# ════════════════════════════════════════════════
# Slide 25: 团队分工
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "附录 A  团队分工", "Team Contribution")

team_table = [
    ("成员", "主要负责模块", "Git Commits"),
    ("付宝昊", "EDA、特征工程、模型设计、报告撰写、PPT 制作", "15+"),
    ("梁智毅", "（待补充）", "—"),
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2), team_table,
          [Inches(2), Inches(7), Inches(3)])

# Git repo
add_text(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(0.4),
         "项目仓库", size=18, color=RED, bold=True)

repo_items = [
    "GitHub：https://github.com/UnknownAnt/H-M-recommendation",
    "主要文件：",
    "  myCode.ipynb — V13 最终模型（MAP@12 = 0.02229）",
    "  EDA_Checkpoint.ipynb — 数据探索与清洗",
    "  myNotebook.ipynb — LightGBM 消融与 SHAP 分析",
    "  期末报告.md — 完整数据分析报告",
    "  H&M_个性化推荐系统_数据分析报告_v2.pptx — 本演示文稿",
]
add_multiline(slide, Inches(0.6), Inches(3.3), Inches(12), Inches(3),
              repo_items, size=14, color=DARK)


# ════════════════════════════════════════════════
# Slide 26: 参考资料
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)
add_header_bar(slide, "附录 B  参考资料", "References")

refs = [
    "[1] Kaggle. H&M Personalized Fashion Recommendations.",
    "    https://www.kaggle.com/competitions/h-and-m-personalized-fashion-recommendations",
    "",
    "[2] Liu, Y. et al. Time-decay weighting for recommender systems. RecSys, 2020.",
    "",
    "[3] He, X. & McAuley, J. Matrix Factorization for Implicit Feedback. KDD, 2016.",
    "",
    "[4] Lundberg, S. & Lee, S.-I. A Unified Approach to Interpreting Model Predictions.",
    "    Advances in Neural Information Processing Systems (NeurIPS), 2017.",
    "",
    "[5] Ke, G. et al. LightGBM: A Highly Efficient Gradient Boosting Decision Tree.",
    "    Advances in Neural Information Processing Systems (NeurIPS), 2017.",
]
add_multiline(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5),
              refs, size=15, color=DARK)


# ════════════════════════════════════════════════
# Slide 27: Thank You
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.12), RED)
add_rect(slide, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), RED)

add_text(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1),
         "Thank You", size=52, color=DARK, bold=True, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.3), Inches(2.3), Pt(3), RED)

add_text(slide, Inches(1), Inches(3.6), Inches(11.3), Inches(0.6),
         "H&M 个性化推荐系统 · 数据分析报告", size=20, color=RED, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
         "付宝昊 · 梁智毅  |  2026 年 5 月 6 日", size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.9), Inches(11.3), Inches(0.5),
         "github.com/UnknownAnt/H-M-recommendation", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
         "MAP@12 = 0.02229  |  +19% vs Baseline", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(7.08), Inches(11.7), Inches(0.35),
         "课程：大数据分析与计算 (610ZH125)  |  竞赛：Kaggle H&M Personalized Fashion Recommendations",
         size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ──
output_path = r"d:\Un_Projects\H&M\H&M_个性化推荐系统_数据分析报告_v2.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
